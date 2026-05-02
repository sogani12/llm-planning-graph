"""
Design constants and python-pptx helper functions.
Dark navy theme with sky-blue headers and amber accents.
"""
from __future__ import annotations
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
BG          = RGBColor(0x0F, 0x17, 0x29)   # deep navy
HEADER      = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
SUBHEADER   = RGBColor(0xF5, 0x9E, 0x0B)   # amber
BODY        = RGBColor(0xE2, 0xE8, 0xF0)   # off-white
DIM         = RGBColor(0x94, 0xA3, 0xB8)   # muted slate
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# Chart colours (also used in matplotlib)
C_BLUE   = "#38BDF8"
C_AMBER  = "#F59E0B"
C_GREEN  = "#10B981"
C_GRAY   = "#64748B"
C_BG     = "#0F1729"
C_TEXT   = "#E2E8F0"
C_DIM    = "#94A3B8"

# ---------------------------------------------------------------------------
# Slide dimensions  (16:9 widescreen, 10 × 5.625 inches)
# ---------------------------------------------------------------------------
W = Inches(10)
H = Inches(5.625)

MARGIN_L = Inches(0.45)
MARGIN_T = Inches(0.30)
CONTENT_W = Inches(9.10)

HEADER_H  = Inches(0.65)
CONTENT_T = Inches(1.05)
CONTENT_H = Inches(4.35)


# ---------------------------------------------------------------------------
# Presentation factory
# ---------------------------------------------------------------------------
def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def blank_slide(prs: Presentation):
    """Add a blank slide (no layout placeholders)."""
    blank_layout = prs.slide_layouts[6]  # index 6 = blank
    return prs.slides.add_slide(blank_layout)


def set_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, text: str,
               top: Emu = MARGIN_T,
               left: Emu = MARGIN_L,
               width: Emu = CONTENT_W,
               height: Emu = HEADER_H,
               font_size: int = 34) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = HEADER


def add_rule(slide,
             top: Emu | None = None,
             left: Emu = MARGIN_L,
             width: Emu = CONTENT_W,
             thickness: Emu = Inches(0.02),
             color: RGBColor = SUBHEADER) -> None:
    """Thin horizontal rule below the header."""
    t = top if top is not None else Emu(int(MARGIN_T) + int(HEADER_H) - Inches(0.05))
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, t, width, thickness,
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_label(slide, text: str, left: Emu, top: Emu,
              width: Emu = Inches(4), height: Emu = Inches(0.4),
              font_size: int = 17) -> None:
    """Amber section label."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = SUBHEADER


def add_body(slide, bullets: list[str | tuple],
             left: Emu = MARGIN_L,
             top: Emu = CONTENT_T,
             width: Emu = CONTENT_W,
             height: Emu = CONTENT_H,
             font_size: int = 16,
             line_spacing: float = 1.25) -> None:
    """
    Add a bullet list. Each item can be:
      str            → top-level bullet
      (str, int)     → (text, indent_level)  level 0 = top, 1 = sub
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = level
        # line spacing
        p.space_before = Pt(3 if level == 0 else 1)

        # bullet character
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•' if level == 0 else '◦')

        # indent
        pPr.set('indent', str(int(Inches(-0.25))))
        pPr.set('marL', str(int(Inches(0.30 + level * 0.25))))

        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size - level * 1.5)
        run.font.color.rgb = BODY if level == 0 else DIM
        run.font.bold = False

        # bullet colour (amber)
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', 'F59E0B')

        buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
        buSzPct.set('val', '90000')


def embed_image(slide, buf: BytesIO,
                left: Emu, top: Emu,
                width: Emu, height: Emu) -> None:
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width, height)


def add_text_box(slide, text: str,
                 left: Emu, top: Emu, width: Emu, height: Emu,
                 font_size: int = 15,
                 color: RGBColor = BODY,
                 bold: bool = False,
                 align=PP_ALIGN.LEFT) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
