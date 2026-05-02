"""
Build Team-16-v2.pptx — dark-theme deck with EDA charts.

Usage:
    python ppt/scripts/build_slides.py
Output:
    ppt/Team-16-v2.pptx
"""
from __future__ import annotations
import sys
from pathlib import Path

# Make theme/charts importable when run from any cwd
sys.path.insert(0, str(Path(__file__).parent))

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import charts
from theme import (
    new_presentation, blank_slide, set_bg,
    add_header, add_rule, add_label, add_body, embed_image, add_text_box,
    BG, HEADER, SUBHEADER, BODY, DIM, WHITE,
    MARGIN_L, MARGIN_T, CONTENT_W, CONTENT_T, CONTENT_H, HEADER_H,
    W, H,
)
from pptx.util import Inches, Emu

OUT = Path(__file__).parent.parent / "Team-16-v2.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def col2(slide, left_bullets, right_bullets,
         left_label="", right_label="",
         split=0.50, top=CONTENT_T, font_size=15):
    """Two-column bullet layout."""
    lw = Inches(9.10 * split - 0.15)
    rw = Inches(9.10 * (1 - split) - 0.15)
    rl = MARGIN_L + lw + Inches(0.30)

    # Labels sit between the rule (~0.90") and the bullets — fixed gap below rule
    label_top = Inches(0.97)
    if left_label:
        add_label(slide, left_label, MARGIN_L, label_top, lw)
    if right_label:
        add_label(slide, right_label, rl, label_top, rw)

    # Body starts below label (label height 0.40" + 0.08" gap = 0.48" below label_top)
    has_labels = bool(left_label or right_label)
    body_top = label_top + Inches(0.48) if has_labels else top

    if left_bullets:
        add_body(slide, left_bullets, MARGIN_L, body_top, lw,
                 CONTENT_H - Inches(0.4), font_size)
    if right_bullets:
        add_body(slide, right_bullets, rl, body_top, rw,
                 CONTENT_H - Inches(0.4), font_size)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def s01_title(prs):
    slide = blank_slide(prs)
    set_bg(slide)

    # Centre block
    add_text_box(slide,
        "Structured Planning Under Uncertainty:",
        Inches(0.5), Inches(1.3), Inches(9), Inches(0.55),
        font_size=28, color=HEADER, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide,
        "A Decision Graph Approach to AI-Assisted Software Development",
        Inches(0.5), Inches(1.85), Inches(9), Inches(0.55),
        font_size=22, color=HEADER, bold=True, align=PP_ALIGN.CENTER)

    add_rule(slide, top=Inches(2.52),
             left=Inches(2.0), width=Inches(6.0))

    add_text_box(slide,
        "Naman Sogani · Ashvin Sehgal · Mari Garey · Elena Motonishi · Sreemae Konduru · Omkar Khade",
        Inches(0.5), Inches(2.70), Inches(9), Inches(0.4),
        font_size=14, color=BODY, align=PP_ALIGN.CENTER)
    add_text_box(slide,
        "CS 639 — Deep Learning for NLP  ·  UW–Madison  ·  Spring 2026",
        Inches(0.5), Inches(3.15), Inches(9), Inches(0.4),
        font_size=13, color=DIM, align=PP_ALIGN.CENTER)


def s02_problem(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Problem & Motivation")
    add_rule(slide)

    col2(slide,
        left_label="The Gap",
        left_bullets=[
            "AI coding tools optimise for the immediate task — they have no memory of past decisions, system-wide constraints, or architectural intent",
            ("Past decisions and coding standards are silently ignored", 1),
            ("Cross-module dependencies are routinely missed", 1),
            ("Errors surface later in unrelated parts of the codebase", 1),
        ],
        right_label="Why It Matters",
        right_bullets=[
            "As projects grow, inconsistency compounds — a single misaligned decision can invalidate downstream components",
            "No existing tool surfaces the planning graph implicit in developer–AI conversations",
            "We propose extracting that graph and using it to keep AI grounded in project-wide context",
        ],
        top=CONTENT_T + Inches(0.15),
        font_size=15,
    )


def s03_related(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Related Work")
    add_rule(slide)

    papers = [
        ("Graph Chain-of-Thought  (Jin et al., 2024)",
         "Graph-structured reasoning over multi-step dependencies via nodes & edges — but no persistent or evolving system context."),
        ("Search-R1  (Jin et al., 2025)",
         "RL-trained LLM that decides when to search for external info — improves reasoning via adaptive retrieval, but works on unstructured data only."),
        ("GNN-RAG  (Mavromatis & Karypis, 2025)",
         "GNNs retrieve relevant subgraphs for multi-hop LLM reasoning — but relies on static, pre-built graphs and non-adaptive retrieval."),
    ]

    top = CONTENT_T + Inches(0.05)
    for title, body in papers:
        add_label(slide, f"▸  {title}", MARGIN_L, top, CONTENT_W, Inches(0.35), font_size=16)
        add_body(slide, [("→  " + body, 1)],
                 MARGIN_L + Inches(0.20), top + Inches(0.33),
                 CONTENT_W - Inches(0.20), Inches(0.55), font_size=14)
        top += Inches(1.10)


def s04_approach(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Our Approach: Decision Graph Extraction")
    add_rule(slide)

    col2(slide,
        left_label="7 Node Types",
        left_bullets=[
            "Objective   — what the system aims to achieve",
            "Requirement — hard constraints and must-haves",
            "Assumption  — implicit beliefs about the environment",
            "Decision    — explicit architectural choices + rationale",
            "Component   — modules, services, libraries",
            "Interface   — APIs, protocols, schemas",
            "Risk        — concerns, tradeoffs, failure modes",
        ],
        right_label="8 Edge Types",
        right_bullets=[
            "motivated_by    — decision ← objective",
            "assumes         — decision → assumption",
            "implements      — component → requirement",
            "depends_on      — component → component",
            "conflicts_with  — risk ↔ decision",
            "invalidates     — risk → assumption",
            "exposes         — component → interface",
            "consumes        — component → interface",
        ],
        top=CONTENT_T + Inches(0.10),
        font_size=14,
    )

    add_text_box(slide,
        "Extracted from developer–AI conversation logs via LLM; stored as a typed directed graph (NetworkX).",
        MARGIN_L, H - Inches(0.65), CONTENT_W, Inches(0.40),
        font_size=13, color=DIM)


def s05_eda_method(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Data Collection & EDA Methodology")
    add_rule(slide)

    col2(slide,
        left_label="Corpus (166 documents, 5,702 sentences)",
        left_bullets=[
            "GitHub Issues — 116 issues across 8 design-heavy repos fetched via Search API (keyword: design OR architecture OR rfc OR proposal OR decision OR tradeoff in title)",
            ("rust-lang/rfcs, denoland/deno, astral-sh/uv, tokio-rs/tokio, gleam-lang/gleam, ghostty-org/ghostty, BurntSushi/ripgrep, microsoft/typescript", 1),
            "Stack Overflow — 30 [software-design] questions + top 3 voted answers per question",
            "Post-mortems — danluu/post-mortems curated list (8k words of incident descriptions)",
        ],
        right_label="4-Phase EDA",
        right_bullets=[
            "Hedge detection — epistemic modals & conditionals as Assumption node signal",
            "Node-type keyword density — hits per 1k words across 7 node types × 3 tiers",
            "Dependency verb patterns — surface signals for edge extraction",
            "Uncategorized sampling — manual review to surface schema gaps",
        ],
        top=CONTENT_T + Inches(0.10),
        font_size=14,
    )


def s06_eda_density(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "EDA Findings: Node-Type Keyword Density")
    add_rule(slide)

    buf = charts.node_density_chart()
    embed_image(slide, buf,
                MARGIN_L, CONTENT_T - Inches(0.05),
                Inches(9.10), Inches(4.10))

    add_text_box(slide,
        "Assumption is rarest across all tiers (mean 0.049/1k).  "
        "GitHub Issues are the only reliable source of Risk & Decision signal.",
        MARGIN_L, H - Inches(0.50), CONTENT_W, Inches(0.38),
        font_size=12, color=DIM)


def s07_eda_signal(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "EDA Findings: Hedge & Edge Signal")
    add_rule(slide)

    # Two charts side by side
    buf_h = charts.hedge_chart()
    buf_d = charts.dependency_chart()

    embed_image(slide, buf_h,
                MARGIN_L, CONTENT_T,
                Inches(4.45), Inches(3.85))
    embed_image(slide, buf_d,
                MARGIN_L + Inches(4.65), CONTENT_T,
                Inches(4.45), Inches(3.85))

    add_text_box(slide,
        "320 hedge hits across 5,702 sentences.  "
        "11/14 dependency verb patterns found; 'consumes', 'feeds_into', 'delegates_to' absent from corpus.",
        MARGIN_L, H - Inches(0.50), CONTENT_W, Inches(0.38),
        font_size=12, color=DIM)


def s08_prefix(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Prefix Tuning: Repo-Aware LLM Specialisation")
    add_rule(slide)

    col2(slide,
        left_label="What It Does",
        left_bullets=[
            "Trains separate PEFT prefix adapters (20 virtual tokens) per repo domain — e.g. finance_python, backend_sql",
            "JSON metadata router selects the right adapter at inference time based on frameworks, directive, and file-path constraints",
            "Soft-biases LLM output style without storing repo data in weights — retrieval handles knowledge, prefix handles behaviour",
        ],
        right_label="Why It Fits",
        right_bullets=[
            "GitHub Issues (Decision/Risk-heavy) and Stack Overflow (Component/Interface-heavy) have near-disjoint node profiles → natural adapter split",
            "Modular: new domains added as small adapters, not new models",
            "Next: build train_examples.json and run first adapter training loop",
        ],
        top=CONTENT_T + Inches(0.10),
        font_size=15,
    )


def s09_eval(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Evaluation Plan")
    add_rule(slide)

    col2(slide,
        left_label="Metrics",
        left_bullets=[
            "Node Recall — fraction of ground-truth nodes recovered from a conversation",
            "Edge Precision — fraction of extracted edges that are correct",
            "Failure Point Recall — fraction of post-mortem failure points mapped to Risk or Assumption nodes",
        ],
        right_label="Experiment Design",
        right_bullets=[
            "6 planning projects spanning diverse complexity (job pipeline, OAuth2/SSO, ETL, monolith migration, ML feature, dev CLI)",
            "Condition A: standard Claude Code session — prompt in, plan out",
            ("Condition B: graph-augmented — extraction after each exchange, graph fed back as planning context", 1),
            "Each project scored against a pre-defined objectives rubric (functional + structural)",
            "Baseline: RAG-augmented planning to isolate graph contribution from retrieval alone",
        ],
        top=CONTENT_T + Inches(0.10),
        font_size=14,
    )

    add_label(slide, "Dataset", MARGIN_L, H - Inches(1.15), CONTENT_W, Inches(0.35))
    add_body(slide,
        ["Open-sourced codebases with documented design decisions (sourced from our EDA corpus); "
         "manually annotated decision graphs as ground truth."],
        MARGIN_L, H - Inches(0.72), CONTENT_W, Inches(0.50),
        font_size=14)


def s10_findings(prs):
    from pptx.dml.color import RGBColor as _RGB

    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Case Study: Job Copilot Pipeline")
    add_rule(slide)

    RED     = _RGB(0xF8, 0x71, 0x71)
    GREEN   = _RGB(0x10, 0xB9, 0x81)
    AMBER_C = _RGB(0xF5, 0x9E, 0x0B)
    HDR_A   = _RGB(0x2A, 0x1A, 0x0E)
    HDR_B   = _RGB(0x0D, 0x26, 0x1C)
    HDR_L   = _RGB(0x12, 0x1E, 0x32)
    ROW_DARK = _RGB(0x12, 0x1D, 0x30)
    ROW_MID  = _RGB(0x0F, 0x17, 0x29)

    # (label, a_text, a_color, b_text, b_color)
    rows = [
        ("Functional (7)",  "5.5 / 7",                     AMBER_C, "6 / 7",                           AMBER_C),
        ("Structural (3)",  "1 / 3",                        RED,     "3 / 3",                           GREEN),
        ("Runtime",         "Python — implicit ✗",          RED,     "Python — documented\n+ alternatives ✓", AMBER_C),
        ("Assumptions",     "None enumerated",              RED,     "8 nodes,\nconfidence-rated",      GREEN),
        ("Risks",           "6 items, informal prose",      AMBER_C, "9 nodes,\nseverity-rated",        GREEN),
        ("Graph",           "—",                            DIM,     "57 nodes · 64 edges",             BODY),
    ]

    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(
        n_rows, 3,
        MARGIN_L, CONTENT_T,
        CONTENT_W, Inches(3.10),
    ).table

    tbl.columns[0].width = Inches(1.80)
    tbl.columns[1].width = Inches(3.65)
    tbl.columns[2].width = Inches(3.65)

    def _cell(r, c, text, fg=BODY, bg=None, bold=False, fs=13,
              align=PP_ALIGN.CENTER):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg if bg else ROW_MID
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.color.rgb = fg
        run.font.bold = bold

    _cell(0, 0, "",            fg=DIM,     bg=HDR_L, bold=True, fs=14)
    _cell(0, 1, "Condition A", fg=AMBER_C, bg=HDR_A, bold=True, fs=15)
    _cell(0, 2, "Condition B", fg=GREEN,   bg=HDR_B, bold=True, fs=15)

    for i, (label, a_txt, a_col, b_txt, b_col) in enumerate(rows):
        bg = ROW_DARK if i % 2 == 0 else ROW_MID
        fs = 15 if i < 2 else 13
        _cell(i + 1, 0, label, fg=DIM,   bg=bg, bold=True, fs=13, align=PP_ALIGN.LEFT)
        _cell(i + 1, 1, a_txt, fg=a_col, bg=bg, bold=(i < 2), fs=fs)
        _cell(i + 1, 2, b_txt, fg=b_col, bg=bg, bold=(i < 2), fs=fs)

    # Key findings below the table
    findings_top = CONTENT_T + Inches(3.18)
    add_label(slide, "Key Findings", MARGIN_L, findings_top, CONTENT_W, Inches(0.30), font_size=14)
    add_body(slide, [
        "Structural coverage dominated even with a broken feedback loop — B surfaced 8 explicit assumption nodes + 9 risk nodes; A produced zero enumerable ones",
        "Both chose Python for Playwright (wrong); B's graph logged the decision with alternatives_considered — the mistake is auditable and challengeable",
        "Core result: the graph forces explicit enumeration of assumptions and risks regardless of whether the planner makes correct decisions",
    ],
    MARGIN_L, findings_top + Inches(0.33), CONTENT_W, Inches(1.10),
    font_size=11)


def s11_graph(prs):
    import os
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Decision Graph: Job Copilot Pipeline (Condition B)")
    add_rule(slide)

    img_path = str(Path(__file__).parent.parent / "assets" / "graph-proj1-v2.png")
    slide.shapes.add_picture(
        img_path,
        MARGIN_L, CONTENT_T,
        Inches(7.60), Inches(3.85),
    )

    add_text_box(slide,
        "57 nodes · 64 edges · extracted from 2 planning exchanges. "
        "Colours: blue = objective, orange = requirement, green = decision, "
        "teal = component, purple = interface, yellow = assumption, red = risk.",
        MARGIN_L + Inches(7.70), CONTENT_T,
        Inches(1.75), Inches(3.85),
        font_size=11, color=DIM)


def s12_next(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header(slide, "Next Steps")
    add_rule(slide)

    add_label(slide, "Immediate Priorities", MARGIN_L, Inches(0.97), CONTENT_W)
    add_body(slide, [
        "Run comparative evaluation on 5 remaining projects — vary domain (OAuth2/SSO, ETL pipeline, monolith migration, ML feature, dev CLI) to test graph generalization",
        "Test true iterative feedback loop: two-window workflow where graph is re-injected after each decision cluster — isolate whether iteration improves planning vs. single-shot extraction",
        "Experiment with stronger Condition B prompting mechanisms: 'reference graph before each decision', structured decision templates, explicit assumption-surfacing prompts",
        "Compare planning agents: Claude Code vs Cursor vs raw API — do different agents respond differently to graph context?",
        "Annotate ground truth graphs + train prefix-tuning adapters for domain-specific assumption and risk extraction",
    ],
    MARGIN_L, Inches(1.45), CONTENT_W, CONTENT_H - Inches(0.5),
    font_size=15)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = new_presentation()

    s01_title(prs)
    s02_problem(prs)
    s03_related(prs)
    s04_approach(prs)
    s05_eda_method(prs)
    s06_eda_density(prs)
    s07_eda_signal(prs)
    s08_prefix(prs)
    s09_eval(prs)
    s10_findings(prs)
    s11_graph(prs)
    s12_next(prs)

    prs.save(OUT)
    print(f"Saved → {OUT}")


if __name__ == "__main__":
    main()
